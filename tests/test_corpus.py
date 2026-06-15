from pathlib import Path

from docx import Document as DocxDocument
from genacademy_rag.core.types import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from genacademy_coach.corpus import (
    build_doc_id,
    extraction_summary,
    load_docx_document,
    load_markdown_document,
    load_pptx_document,
    load_pptx_document_with_stats,
    source_type_for_path,
)


def test_source_type_for_path_uses_folder_priority():
    assert source_type_for_path(Path("corpus/slides/week1-session1.pptx")) == "slide"
    assert source_type_for_path(Path("corpus/handouts/agent-memory.pdf")) == "handout"
    assert source_type_for_path(Path("corpus/notes/lesson1.md")) == "note"
    assert source_type_for_path(Path("corpus/transcripts/week1-session1.md")) == "transcript"


def test_build_doc_id_is_stable_and_does_not_include_absolute_path():
    first = build_doc_id(Path("corpus/notes/lesson1.md"), b"hello")
    second = build_doc_id(Path("/tmp/elsewhere/corpus/notes/lesson1.md"), b"hello")

    assert first == second
    assert first.startswith("note/lesson1-")


def test_load_markdown_document_sets_source_type_and_title(tmp_path):
    path = tmp_path / "corpus" / "notes" / "lesson1.md"
    path.parent.mkdir(parents=True)
    path.write_text("# Attention\n\nText.", encoding="utf-8")

    doc = load_markdown_document(path)

    assert doc.source_type == "note"
    assert doc.title == "lesson1.md"
    assert "Attention" in doc.text


def test_load_docx_document_includes_paragraphs_and_tables(tmp_path):
    path = tmp_path / "corpus" / "handouts" / "lesson.docx"
    path.parent.mkdir(parents=True)
    parsed = DocxDocument()
    parsed.add_paragraph("Paragraph text")
    table = parsed.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Term"
    table.cell(0, 1).text = "Definition"
    parsed.save(path)

    doc = load_docx_document(path)

    assert doc.source_type == "handout"
    assert "Paragraph text" in doc.text
    assert "Term" in doc.text
    assert "Definition" in doc.text


def test_load_pptx_document_includes_shape_text_and_speaker_notes(tmp_path):
    path = tmp_path / "corpus" / "slides" / "week1.pptx"
    path.parent.mkdir(parents=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Retrieval planning"
    slide.notes_slide.notes_text_frame.text = "Speaker note about citations."
    prs.save(path)

    doc = load_pptx_document(path)

    assert doc.source_type == "slide"
    assert "Retrieval planning" in doc.text
    assert "Speaker note about citations" in doc.text


def test_image_only_pptx_reports_empty_text_with_shape_count(tmp_path):
    path = tmp_path / "corpus" / "slides" / "image-only.pptx"
    path.parent.mkdir(parents=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    prs.save(path)

    loaded = load_pptx_document_with_stats(path)
    doc = loaded.document
    summary = extraction_summary(doc, slide_shape_count=loaded.slide_shape_count)

    assert doc.text.strip() == ""
    assert summary["empty"] is True
    assert summary["slide_shape_count"] == 1


def test_pptx_loader_returns_shape_count_from_single_parse(tmp_path):
    path = tmp_path / "corpus" / "slides" / "shape-count.pptx"
    path.parent.mkdir(parents=True)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2), Inches(2), Inches(1))
    prs.save(path)

    loaded = load_pptx_document_with_stats(path)

    assert loaded.document.source_type == "slide"
    assert loaded.slide_shape_count == 2


def test_extraction_summary_marks_empty_text():
    doc = Document(
        doc_id="note/empty",
        title="empty.md",
        source_type="note",
        text="  ",
        stored_path="corpus/notes/empty.md",
    )

    summary = extraction_summary(doc)

    assert summary["chars"] == 0
    assert summary["empty"] is True
