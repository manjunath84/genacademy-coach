from __future__ import annotations

import hashlib
from dataclasses import replace
from pathlib import Path

from docx import Document as DocxDocument
from genacademy_rag.core.loaders.pdf_loader import load_pdf_bytes
from genacademy_rag.core.types import Document
from pptx import Presentation

INDEXABLE_DIRS = ("notes", "transcripts", "slides", "handouts")
INDEXABLE_SUFFIXES = {".md", ".pdf", ".pptx", ".docx"}


def source_type_for_path(path: Path) -> str:
    parts = set(path.parts)
    if "slides" in parts:
        return "slide"
    if "handouts" in parts:
        return "handout"
    if "notes" in parts:
        return "note"
    if "transcripts" in parts:
        return "transcript"
    raise ValueError(f"not an indexable corpus path: {path}")


def build_doc_id(path: Path, raw_bytes: bytes) -> str:
    source_type = source_type_for_path(path)
    stem = path.stem.lower().replace("_", "-")
    digest = hashlib.sha256(raw_bytes).hexdigest()[:12]
    return f"{source_type}/{stem}-{digest}"


def load_markdown_document(path: Path) -> Document:
    raw = path.read_bytes()
    return Document(
        doc_id=build_doc_id(path, raw),
        title=path.name,
        source_type=source_type_for_path(path),
        text=raw.decode("utf-8"),
        filename=path.name,
        stored_path=str(path),
    )


def load_pdf_document(path: Path) -> Document:
    raw = path.read_bytes()
    doc = load_pdf_bytes(filename=path.name, raw_bytes=raw, stored_path=str(path))
    return replace(
        doc,
        doc_id=build_doc_id(path, raw),
        source_type=source_type_for_path(path),
        title=path.name,
    )


def _shape_alt_text(shape) -> list[str]:
    values: list[str] = []
    for node in shape.element.iter():
        if not str(node.tag).endswith("}cNvPr"):
            continue
        for key in ("title", "descr"):
            value = (node.attrib.get(key) or "").strip()
            if value:
                values.append(value)
    return values


def _shape_text(shape) -> list[str]:
    values: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text_frame.text.strip()
        if text:
            values.append(text)
    if getattr(shape, "has_table", False):
        cells = [
            cell.text.strip()
            for row in shape.table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        if cells:
            values.append("\n".join(cells))
    values.extend(_shape_alt_text(shape))
    return values


def _slide_notes_text(slide) -> str:
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except (AttributeError, KeyError, ValueError):
        return ""


def pptx_shape_count(path: Path) -> int:
    prs = Presentation(path)
    return sum(len(slide.shapes) for slide in prs.slides)


def load_pptx_document(path: Path) -> Document:
    raw = path.read_bytes()
    prs = Presentation(path)
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            texts.extend(_shape_text(shape))
        notes = _slide_notes_text(slide)
        if notes:
            texts.append("## Speaker Notes\n\n" + notes)
        if texts:
            parts.append(f"# Slide {idx}\n\n" + "\n\n".join(texts))
    return Document(
        doc_id=build_doc_id(path, raw),
        title=path.name,
        source_type="slide",
        text="\n\f\n".join(parts),
        filename=path.name,
        stored_path=str(path),
    )


def _docx_table_text(parsed) -> list[str]:
    return [
        cell.text.strip()
        for table in parsed.tables
        for row in table.rows
        for cell in row.cells
        if cell.text.strip()
    ]


def load_docx_document(path: Path) -> Document:
    raw = path.read_bytes()
    parsed = DocxDocument(path)
    paragraphs = [p.text.strip() for p in parsed.paragraphs if p.text.strip()]
    table_cells = _docx_table_text(parsed)
    return Document(
        doc_id=build_doc_id(path, raw),
        title=path.name,
        source_type=source_type_for_path(path),
        text="\n\n".join([*paragraphs, *table_cells]),
        filename=path.name,
        stored_path=str(path),
    )


def load_corpus_document(path: Path) -> Document:
    suffix = path.suffix.lower()
    if suffix == ".md":
        return load_markdown_document(path)
    if suffix == ".pdf":
        return load_pdf_document(path)
    if suffix == ".pptx":
        return load_pptx_document(path)
    if suffix == ".docx":
        return load_docx_document(path)
    raise ValueError(f"unsupported corpus file: {path}")


def iter_indexable_files(corpus_dir: Path) -> list[Path]:
    files: list[Path] = []
    for dirname in INDEXABLE_DIRS:
        root = corpus_dir / dirname
        if not root.exists():
            continue
        files.extend(
            path
            for path in root.rglob("*")
            if path.is_file() and path.suffix.lower() in INDEXABLE_SUFFIXES
        )
    return sorted(files)


def extraction_summary(doc: Document) -> dict[str, object]:
    text = doc.text.strip()
    stored_path = Path(doc.stored_path) if doc.stored_path else None
    slide_shape_count = None
    if (
        doc.source_type == "slide"
        and stored_path is not None
        and stored_path.suffix.lower() == ".pptx"
        and stored_path.exists()
    ):
        slide_shape_count = pptx_shape_count(stored_path)
    return {
        "doc_id": doc.doc_id,
        "title": doc.title,
        "source_type": doc.source_type,
        "chars": len(text),
        "empty": not bool(text),
        "slide_shape_count": slide_shape_count,
        "stored_path": doc.stored_path,
    }
